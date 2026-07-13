# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
NAVY   = RGBColor(0x1F, 0x2A, 0x44)
BLUE   = RGBColor(0x2F, 0x6F, 0xED)
TEAL   = RGBColor(0x0E, 0x9E, 0x9E)
CORAL  = RGBColor(0xE8, 0x5D, 0x4E)
AMBER  = RGBColor(0xF2, 0xA6, 0x2E)
GRAY   = RGBColor(0x5B, 0x63, 0x72)
LIGHT  = RGBColor(0xF3, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARKTX = RGBColor(0x25, 0x2B, 0x38)

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, text, size=18, color=DARKTX, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def header(s, kicker, title, tcolor=NAVY):
    rect(s, 0, 0, SW, Inches(1.35), WHITE)
    rect(s, Inches(0.55), Inches(0.42), Inches(0.14), Inches(0.55), BLUE)
    txt(s, Inches(0.85), Inches(0.30), Inches(11.8), Inches(0.35), kicker, 12.5, TEAL, bold=True)
    txt(s, Inches(0.85), Inches(0.58), Inches(11.8), Inches(0.6), title, 26, tcolor, bold=True)
    rect(s, 0, Inches(1.35), SW, Pt(2.2), LIGHT)

def pagenum(s, n):
    txt(s, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35), str(n), 11, GRAY, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- 1. TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(5.05), SW, Inches(0.06), BLUE)
rect(s, Inches(0.9), Inches(1.75), Inches(0.9), Inches(0.13), AMBER)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5), "소셜모닝살롱 · AI 사전설문", 18, RGBColor(0xBF,0xD0,0xF2), bold=True)
txt(s, Inches(0.9), Inches(2.6), Inches(11.8), Inches(1.6),
    "AI 사전설문 응답 분석", 46, WHITE, bold=True)
txt(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.6),
    "2시간 모닝살롱 세션 설계를 위한 참가자 인사이트", 20, RGBColor(0xD7,0xE0,0xF2))
txt(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.5),
    "응답 6명   ·   2026.07.12–07.13   ·   분석일 2026.07.13", 14, RGBColor(0x9F,0xB2,0xDB))

# ---------------------------------------------------------------- 2. SUMMARY
s = slide()
header(s, "EXECUTIVE SUMMARY", "핵심 요약 — 한눈에 보기")
cards = [
    ("이미 쓰는 실무자", "6명 전원 소셜/임팩트 리더급.\n초심자가 아니라 '매일 쓰지만\n얕게 쓰는' 그룹", TEAL),
    ("자동화의 벽", "전원의 공통 페인포인트:\n\"에이전트·자동화는 하고 싶은데\n너무 기술적으로 느껴진다\"", CORAL),
    ("원하는 것", "내 업무에 바로 붙이는\nAI 자동화·에이전트.\n제안서·회의록·리포트 반복업무", BLUE),
    ("원하는 방식", "실제 업무 사례 데모 +\n프롬프트/워크시트 제공.\n개념·입문 강의는 최소화", AMBER),
]
x0 = Inches(0.55); gap = Inches(0.25)
cw = (SW - x0*2 - gap*3) / 4
cy = Inches(1.9); ch = Inches(4.4)
for i,(t,b,c) in enumerate(cards):
    x = x0 + i*(cw+gap)
    rect(s, x, cy, cw, ch, LIGHT)
    rect(s, x, cy, cw, Inches(0.18), c)
    txt(s, x+Inches(0.2), cy+Inches(0.45), cw-Inches(0.4), Inches(1.0), t, 18, c, bold=True)
    txt(s, x+Inches(0.2), cy+Inches(1.5), cw-Inches(0.4), ch-Inches(1.7), b, 13.5, DARKTX, spacing=1.15)
txt(s, x0, Inches(6.5), SW-x0*2, Inches(0.5),
    "→ 세션의 성패는 '비개발자도 만드는 자동화·에이전트'의 진입 장벽을 낮추는 데 달려 있음", 15, NAVY, bold=True)
pagenum(s, 2)

# ---------------------------------------------------------------- 3. PROFILE (table)
s = slide()
header(s, "RESPONDENTS", "응답자 프로필 — 전원 소셜임팩트 섹터")
rows = [
    ["이름","역할","AI 빈도","자기 수준","파트너 자신감"],
    ["김진남","NPO/국제개발/공공/정책","월 1–2회","Level 1","2"],
    ["한영현","창업가/대표/조직 리더","거의 매일","Level 2","3"],
    ["백정연","사회혁신/소셜벤처 운영자","거의 매일","Level 2","3"],
    ["김하연","창업가/대표/조직 리더","거의 매일","Level 3","4"],
    ["이은애","CSR/ESG/임팩트·기업협력","거의 매일","Level 2","5"],
    ["조정원","교육/연구/컨설팅/교수","주 3–4회","Level 1","1"],
]
nrows, ncols = len(rows), len(rows[0])
tx, ty = Inches(0.55), Inches(1.75)
tw = SW - Inches(1.1)
tbl_shape = s.shapes.add_table(nrows, ncols, tx, ty, tw, Inches(3.9))
table = tbl_shape.table
widths = [Inches(1.5), Inches(4.4), Inches(2.1), Inches(2.1), Inches(2.13)]
for j,wd in enumerate(widths): table.columns[j].width = wd
for i in range(nrows):
    table.rows[i].height = Inches(0.56)
    for j in range(ncols):
        cell = table.cell(i,j)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j<2 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = rows[i][j]
        r.font.name = FONT; r.font.size = Pt(13.5)
        if i==0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            r.font.color.rgb = WHITE; r.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i%2 else LIGHT
            r.font.color.rgb = DARKTX
            if j==4:
                r.font.bold = True; r.font.color.rgb = BLUE
txt(s, tx, Inches(6.05), tw, Inches(1.0),
    "· 섹터 동질성이 높아 하나의 사례(제안서·파트너십·ESG 리포트)가 전원에게 유효\n"
    "· 자신감 편차 큼(1~5) — 조/짝 실습 시 수준 페어링 필요 · Level 4~5(자동화·에이전트)는 0명",
    13.5, GRAY, spacing=1.25)
pagenum(s, 3)

# ---------------------------------------------------------------- 4. USAGE
s = slide()
header(s, "CURRENT USAGE", "AI 활용 현황 — 매일 쓰지만 얕게")
# left: frequency
lx = Inches(0.55); lw = Inches(5.9)
rect(s, lx, Inches(1.85), lw, Inches(0.5), NAVY)
txt(s, lx+Inches(0.2), Inches(1.9), lw-Inches(0.4), Inches(0.4), "사용 빈도", 15, WHITE, bold=True)
freq = [("거의 매일 사용",4,BLUE),("주 3–4회",1,TEAL),("월 1–2회",1,GRAY)]
fy = Inches(2.6)
for label,val,c in freq:
    txt(s, lx, fy, Inches(2.6), Inches(0.4), label, 14, DARKTX)
    barmax = Inches(2.7)
    rect(s, lx+Inches(2.7), fy+Inches(0.05), Emu(int(barmax*val/4)), Inches(0.3), c)
    txt(s, lx+Inches(2.7)+Emu(int(barmax*val/4))+Inches(0.05), fy, Inches(0.6), Inches(0.4), f"{val}명", 13, GRAY, bold=True)
    fy += Inches(0.62)
# self level
rect(s, lx, Inches(4.7), lw, Inches(0.5), NAVY)
txt(s, lx+Inches(0.2), Inches(4.75), lw-Inches(0.4), Inches(0.4), "자기 인식 수준", 15, WHITE, bold=True)
lv = [("Level 1  질문·요약·번역 보조","2명"),("Level 2  문서·리서치·기획 활용","3명"),
      ("Level 3  템플릿·파일·워크스페이스","1명"),("Level 4–5  자동화·에이전트·API","0명")]
ly = Inches(5.35)
for label,v in lv:
    c = CORAL if v=="0명" else DARKTX
    txt(s, lx, ly, Inches(4.5), Inches(0.35), label, 13, DARKTX)
    txt(s, lx+Inches(4.6), ly, Inches(1.2), Inches(0.35), v, 13, c, bold=True)
    ly += Inches(0.44)
# right: confidence
rx = Inches(6.8); rw = Inches(5.98)
rect(s, rx, Inches(1.85), rw, Inches(0.5), NAVY)
txt(s, rx+Inches(0.2), Inches(1.9), rw-Inches(0.4), Inches(0.4), "'AI를 업무 파트너로' 자신감 (1–5)", 15, WHITE, bold=True)
vals = [("조정원",1),("김진남",2),("한영현",3),("백정연",3),("김하연",4),("이은애",5)]
by = Inches(2.6)
for name,v in vals:
    txt(s, rx, by, Inches(1.2), Inches(0.35), name, 13, DARKTX)
    barmax = Inches(3.4)
    c = CORAL if v<=2 else (AMBER if v==3 else TEAL)
    rect(s, rx+Inches(1.3), by+Inches(0.03), Emu(int(barmax*v/5)), Inches(0.28), c)
    txt(s, rx+Inches(1.3)+Emu(int(barmax*v/5))+Inches(0.05), by, Inches(0.5), Inches(0.35), str(v), 12.5, GRAY, bold=True)
    by += Inches(0.5)
txt(s, rx, by+Inches(0.1), rw, Inches(0.8), "평균 3.0 · 편차 큼 → 참가자 간 수준 격차 존재", 13.5, NAVY, bold=True)
pagenum(s, 4)

# ---------------------------------------------------------------- 5. TOOLS
s = slide()
header(s, "TOOLS IN USE", "실제 사용 중인 도구 — ChatGPT·Gemini·Claude 표준")
tools = [("Gemini",6),("ChatGPT (유료 Plus)",5),("Claude Pro",4),
         ("Perplexity",1),("Notion AI",1),("Gamma / Tome (발표)",1),
         ("Cursor / Lovable / Replit (코딩)",1),("ChatGPT 무료",1)]
tx0 = Inches(0.7); ty0 = Inches(2.0); rowh = Inches(0.56)
barmax = Inches(6.5)
for i,(name,v) in enumerate(tools):
    y = ty0 + i*rowh
    txt(s, tx0, y, Inches(3.4), Inches(0.4), name, 14, DARKTX, anchor=MSO_ANCHOR.MIDDLE)
    c = BLUE if v>=4 else (TEAL if v>=2 else GRAY)
    w = Emu(int(barmax*v/6))
    rect(s, tx0+Inches(3.5), y+Inches(0.08), w, Inches(0.32), c)
    txt(s, tx0+Inches(3.5)+w+Inches(0.08), y, Inches(0.8), Inches(0.4), f"{v}명", 13, GRAY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx0, Inches(6.55), Inches(12), Inches(0.7),
    "→ 데모는 ChatGPT·Gemini·Claude 3종 안에서 재현하면 대부분 따라올 수 있음  ·  김하연=얼리어답터, 김진남·조정원=상대적 초심자",
    13.5, NAVY, bold=True)
pagenum(s, 5)

# ---------------------------------------------------------------- 6. NEEDS (areas)
s = slide()
header(s, "NEEDS", "AI를 적용하고 싶은 영역")
areas = [("반복업무 자동화",5,CORAL),("보고서/제안서/문서 작성",4,BLUE),
         ("콘텐츠/브랜딩/디자인",3,TEAL),("데이터 정리/분석",3,TEAL),
         ("모금/영업/파트너십 제안",2,GRAY),("리서치/트렌드 파악",1,GRAY),
         ("CSR/ESG/임팩트 측정",1,GRAY),("교육·워크숍 설계",1,GRAY),
         ("커뮤니티/행사 운영",1,GRAY),("개인 생산성/학습",1,GRAY)]
tx0 = Inches(0.7); ty0 = Inches(1.85); rowh = Inches(0.47)
barmax = Inches(6.2)
for i,(name,v,c) in enumerate(areas):
    y = ty0 + i*rowh
    txt(s, tx0, y, Inches(3.9), Inches(0.4), name, 13.5, DARKTX, anchor=MSO_ANCHOR.MIDDLE)
    w = Emu(int(barmax*v/5))
    rect(s, tx0+Inches(4.0), y+Inches(0.06), w, Inches(0.28), c)
    txt(s, tx0+Inches(4.0)+w+Inches(0.08), y, Inches(0.8), Inches(0.4), f"{v}", 12.5, GRAY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.5),
    "→ 1순위 = 반복업무 자동화, 2순위 = 문서/제안서 작성. 세션 핵심 축을 여기에 맞출 것", 14, NAVY, bold=True)
pagenum(s, 6)

# ---------------------------------------------------------------- 7. PAIN POINTS (key)
s = slide()
header(s, "KEY INSIGHT", "가장 크게 막히는 지점 — 세션이 풀어야 할 문제", tcolor=CORAL)
pains = [
    ("자동화·에이전트에 관심 있으나 너무 기술적으로 느껴진다","6명 전원",CORAL,True),
    ("한두 번은 되지만 반복 가능한 워크플로우로 만들기 어렵다","4명",AMBER,True),
    ("도구가 너무 많아 무엇을 써야 할지 모르겠다","2명",GRAY,False),
    ("조직/팀 내 도입 기준을 만들기 어렵다","2명",GRAY,False),
    ("어떤 업무에 써야 할지 모르겠다","1명",GRAY,False),
    ("질문/프롬프트를 어떻게 해야 할지 모르겠다","1명",GRAY,False),
]
y = Inches(1.75)
for name,cnt,c,big in pains:
    h = Inches(0.72) if big else Inches(0.52)
    rect(s, Inches(0.55), y, Inches(9.7), h, LIGHT)
    rect(s, Inches(0.55), y, Inches(0.12), h, c)
    txt(s, Inches(0.85), y, Inches(9.2), h, name, 15 if big else 13.5,
        DARKTX, bold=big, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(10.45), y, Inches(2.3), h, c)
    txt(s, Inches(10.45), y, Inches(2.3), h, cnt, 17 if big else 14, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += h + Inches(0.13)
txt(s, Inches(0.55), y+Inches(0.02), Inches(12.2), Inches(0.5),
    "→ 전원이 \"에이전트/자동화는 하고 싶은데 기술 장벽이 높다\"고 답함. 진입 장벽 낮추기가 핵심.",
    14, NAVY, bold=True)
pagenum(s, 7)

# ---------------------------------------------------------------- 8. EXPECTATION
s = slide()
header(s, "EXPECTATIONS", "세션 기대 — '에이전트'와 '자동화'로 수렴")
quotes = [
    ("김진남","AI agent를 만들고 실제로 활용할 수 있으면"),
    ("한영현","에이전트 만드는 개념이 잡혔으면"),
    ("백정연","AI로 무엇을 시도할지 결정할 수 있음"),
    ("김하연","생산성 극대화 사례를 조직에 접목 · 반복 워크플로우 자동화 · 맞춤형 개인비서 에이전트"),
    ("이은애","2시간의 내용을 잘 소화하면 제일 좋겠다"),
    ("조정원","AI 업무자동화에 관심이 있습니다"),
]
x0=Inches(0.55); gap=Inches(0.3); cw=(SW-x0*2-gap)/2
cy=Inches(1.85); chh=Inches(1.45); gy=Inches(0.2)
for i,(name,q) in enumerate(quotes):
    col=i%2; row=i//2
    x=x0+col*(cw+gap); y=cy+row*(chh+gy)
    rect(s, x, y, cw, chh, LIGHT)
    rect(s, x, y, cw, Inches(0.1), BLUE)
    txt(s, x+Inches(0.25), y+Inches(0.22), cw-Inches(0.5), Inches(0.4), name, 15, BLUE, bold=True)
    txt(s, x+Inches(0.25), y+Inches(0.62), cw-Inches(0.5), chh-Inches(0.7),
        "“"+q+"”", 13.5, DARKTX, italic=True, spacing=1.1)
txt(s, x0, Inches(6.95), SW-x0*2, Inches(0.4),
    "→ 6명 중 최소 4명이 '에이전트/자동화'를 명시적으로 언급", 14, NAVY, bold=True)
pagenum(s, 8)

# ---------------------------------------------------------------- 9. CONTENT PREF
s = slide()
header(s, "CONTENT PREFERENCE", "콘텐츠 선호도 — 무엇을 넣고 무엇을 뺄까")
# want (left)
lx=Inches(0.55); lw=Inches(5.9)
rect(s, lx, Inches(1.8), lw, Inches(0.55), TEAL)
txt(s, lx+Inches(0.2), Inches(1.86), lw-Inches(0.4), Inches(0.45), "✓  원하는 진행 방식", 16, WHITE, bold=True)
want=[("실제 업무 사례 중심 데모","6명"),("프롬프트/워크시트 제공","6명"),
      ("짧고 선명한 인사이트 강의","4명"),("각자 업무 핸즈온 실습","1명")]
y=Inches(2.6)
for a,b in want:
    txt(s, lx+Inches(0.1), y, Inches(4.3), Inches(0.4), "· "+a, 14, DARKTX, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx+Inches(4.5), y, Inches(1.2), Inches(0.4), b, 14, TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.6)
# avoid (right)
rx=Inches(6.85); rw=Inches(5.9)
rect(s, rx, Inches(1.8), rw, Inches(0.55), CORAL)
txt(s, rx+Inches(0.2), Inches(1.86), rw-Inches(0.4), Inches(0.45), "✕  덜 다뤄도 되는 내용", 16, WHITE, bold=True)
avoid=[("AI 역사/개념 설명","4명"),("기초 가입/로그인 안내","4명"),
       ("툴 목록만 나열","3명"),("추상적 미래 전망","2명")]
y=Inches(2.6)
for a,b in avoid:
    txt(s, rx+Inches(0.1), y, Inches(4.3), Inches(0.4), "· "+a, 14, DARKTX, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, rx+Inches(4.5), y, Inches(1.2), Inches(0.4), b, 14, CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.6)
rect(s, Inches(0.55), Inches(5.35), Inches(12.23), Inches(1.45), LIGHT)
rect(s, Inches(0.55), Inches(5.35), Inches(0.12), Inches(1.45), AMBER)
txt(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.2),
    "주의 — '너무 기술적인 코딩 설명' 제외는 단 1명(Level 1)만 요청. 나머지는 오히려 자동화·에이전트 기술을 원함.\n"
    "→ 코딩·기술을 빼면 안 되고, '노코드/로우코드' 비개발자 눈높이로 풀어야 함.",
    14, DARKTX, spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 9)

# ---------------------------------------------------------------- 10. LOGISTICS
s = slide()
header(s, "LOGISTICS", "실습 환경 — 실습 강제 시 2명 소외 위험")
box = [("노트북/태블릿 + AI 계정 모두 준비","4명",TEAL),
       ("스마트폰 + AI 계정만 준비","2명 (김진남·조정원)",AMBER)]
y=Inches(2.0)
for a,b,c in box:
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.95), LIGHT)
    rect(s, Inches(0.7), y, Inches(0.14), Inches(0.95), c)
    txt(s, Inches(1.05), y, Inches(7.5), Inches(0.95), a, 17, DARKTX, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.6), y, Inches(3.9), Inches(0.95), b, 16, c, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.15)
txt(s, Inches(0.7), Inches(4.55), Inches(12), Inches(0.5),
    "· 준비된 4명 전원 PC 사용  ·  실습 성향: '듣고 싶다' 3명 / '실습 준비됨' 3명 (반반)", 14, GRAY)
rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.5), NAVY)
txt(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.15),
    "권고 — 데모 중심 + '원하는 사람은 즉석 실습'의 하이브리드.\n"
    "스마트폰만 가능한 2명을 위해 모바일에서도 따라 할 수 있는 예시 1개를 준비할 것.",
    16, WHITE, bold=True, spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 10)

# ---------------------------------------------------------------- 11. RECOMMENDATIONS
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0.9), Inches(0.7), Inches(0.9), Inches(0.13), AMBER)
txt(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5), "RECOMMENDATIONS", 14, TEAL, bold=True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.7), "세션 설계 제언", 30, WHITE, bold=True)
recs = [
    ("1","핵심 주제 = '비개발자를 위한 AI 자동화·에이전트'","GPTs·Projects·노코드 자동화로 '코딩 없이 만드는 나만의 에이전트'를 보여줄 것"),
    ("2","소재는 이 그룹의 실제 업무로","제안서·파트너십·ESG/임팩트 리포트·회의록 — 하나의 케이스가 전원에게 유효"),
    ("3","짧은 인사이트 강의 → 실무 데모(중심) → 프롬프트/워크시트 배포","개념·역사·가입 안내는 생략, 바로 쓸 프롬프트 팩을 산출물로 제공(전원 기대)"),
    ("4","실습은 '선택형 하이브리드'","준비된 사람은 즉석 실습, 나머지는 데모 관찰 · 모바일 예시 1개 확보"),
    ("5","수준 격차 관리","얼리어답터엔 심화 팁, 초심자엔 '첫 자동화' — 양쪽 모두 가져갈 것 설계"),
]
y=Inches(2.25)
for n,t,d in recs:
    rect(s, Inches(0.9), y, Inches(0.55), Inches(0.55), BLUE)
    txt(s, Inches(0.9), y, Inches(0.55), Inches(0.55), n, 20, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.65), y-Inches(0.03), Inches(11.0), Inches(0.4), t, 16.5, WHITE, bold=True)
    txt(s, Inches(1.65), y+Inches(0.36), Inches(11.0), Inches(0.5), d, 13, RGBColor(0xC7,0xD4,0xEC))
    y+=Inches(0.92)

prs.save("/home/user/003/소셜모닝살롱_AI사전설문_분석.pptx")
print("saved", len(prs.slides.__iter__().__length_hint__() if hasattr(prs.slides,'__length_hint__') else 0))
print("slides:", len(prs.slides._sldIdLst))
